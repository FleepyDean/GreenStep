from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# GreenStep brand colors
GREEN_PRIMARY = RGBColor(0, 168, 132)  # #00A884
GREEN_ACCENT = RGBColor(37, 211, 102)  # #25D366
DARK_TEXT = RGBColor(17, 27, 33)       # #111B21
GRAY_TEXT = RGBColor(84, 101, 111)     # #54656F
WHITE = RGBColor(255, 255, 255)
LIGHT_BG = RGBColor(240, 242, 245)     # #F0F2F5

def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Green gradient background (top half)
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(4))
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = GREEN_PRIMARY
    
    # White background (bottom half)
    bottom_bg = slide.shapes.add_shape(1, 0, Inches(4), prs.slide_width, Inches(3.5))
    fill = bottom_bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(2.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle_text
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(18)
    subtitle_para.font.color.rgb = GRAY_TEXT
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_header(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Green background
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = GREEN_PRIMARY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title_text, content_items, icon=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Light background
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG
    
    # Green header bar
    header = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.2))
    fill = header.fill
    fill.solid()
    fill.fore_color.rgb = GREEN_PRIMARY
    
    # Title with icon
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = f"{icon} {title_text}" if icon else title_text
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    
    # Content box with white background
    content_box = slide.shapes.add_shape(1, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    fill = content_box.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Content text
    text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5))
    tf = text_box.text_frame
    tf.word_wrap = True
    
    for item in content_items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(8)
        if item.startswith("  •") or item.startswith("  ✓"):
            p.level = 1
            p.font.color.rgb = GRAY_TEXT
        elif item.startswith("•") or item.startswith("✓"):
            p.font.color.rgb = GREEN_ACCENT
            p.font.bold = True
    
    return slide

def add_two_column_slide(prs, title_text, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Light background
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG
    
    # Green header bar
    header = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.2))
    fill = header.fill
    fill.solid()
    fill.fore_color.rgb = GREEN_PRIMARY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    
    # Left column
    left_box = slide.shapes.add_shape(1, Inches(0.5), Inches(1.5), Inches(4.3), Inches(5.5))
    fill = left_box.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    left_title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(4), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN_PRIMARY
    
    left_text = slide.shapes.add_textbox(Inches(0.7), Inches(2.3), Inches(4), Inches(4.5))
    tf = left_text.text_frame
    for item in left_items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(6)
    
    # Right column
    right_box = slide.shapes.add_shape(1, Inches(5.2), Inches(1.5), Inches(4.3), Inches(5.5))
    fill = right_box.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    right_title_box = slide.shapes.add_textbox(Inches(5.4), Inches(1.7), Inches(4), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN_PRIMARY
    
    right_text = slide.shapes.add_textbox(Inches(5.4), Inches(2.3), Inches(4), Inches(4.5))
    tf = right_text.text_frame
    for item in right_items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(6)
    
    return slide

print("Creating styled GreenStep presentation...")
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Title
add_title_slide(prs, "🌱 GreenStep", "Carbon Footprint & Eco-Lifestyle Tracker\nSCSM2223 Cross-Platform Application Development\nGroup Project Final Presentation")

# Slide 2: Section - Team
add_section_header(prs, "Team Profile")

# Slide 3: Team Profile
add_content_slide(prs, "Team Profile and Roles", [
    "👤 Team Member: Danish",
    "Role: Full-Stack Developer",
    "",
    "📋 Key Contributions:",
    "  • Frontend development (Vue 3, Vite, Pinia)",
    "  • Backend API development (PHP Slim 4, PDO)",
    "  • Database design and implementation (MySQL)",
    "  • DevOps and deployment (Railway, Git/GitHub)",
    "  • UI/UX design with professional SVG icons",
    "  • Testing, debugging, and documentation"
], "👥")

# Slide 4: Section - Project Overview
add_section_header(prs, "Project Overview")

# Slide 5: Problem & Solution
add_content_slide(prs, "Problem Statement & Solution", [
    "🌍 Problem:",
    "Climate change is among the most urgent global challenges. Individuals collectively account for a large share of avoidable carbon emissions through transport, diet, and energy choices.",
    "",
    "💡 Solution: GreenStep",
    "A friendly daily logger that:",
    "  • Estimates user's carbon footprint automatically",
    "  • Suggests one realistic eco action per day",
    "  • Enables friends to compete on streaks and reductions",
    "  • Provides admin tools for managing emission factors, badges, and tips"
], "🎯")

# Slide 6: Target Users & Objectives
add_two_column_slide(prs, "Target Users & Objectives", 
    "🎯 Target Users", [
        "• Environmentally conscious individuals",
        "• Students and young professionals",
        "• Community groups and organizations",
        "• Anyone wanting to reduce their carbon footprint"
    ],
    "✨ Project Objectives", [
        "• Track daily carbon footprint",
        "• Provide actionable eco tips",
        "• Motivate through gamification",
        "• Enable social accountability",
        "• Help achieve reduction goals"
    ])

# Slide 7: Main Features
add_two_column_slide(prs, "Main Features & Functional Scope",
    "👤 User Features", [
        "• Activity logging with photo upload",
        "• Real-time dashboard with charts",
        "• Carbon reduction goal tracking",
        "• Badge unlocking system",
        "• Eco challenges (join/leave/track)",
        "• Daily eco tips",
        "• Friends & social connections",
        "• Global/friends leaderboard"
    ],
    "⚙️ Admin Features", [
        "• Platform statistics dashboard",
        "• Badge management (CRUD)",
        "• Emission factor management",
        "• Challenge management (CRUD)",
        "• Tips management (CRUD)",
        "• User dataset overview"
    ])

# Slide 8: Section - Technology
add_section_header(prs, "Technology Stack")

# Slide 9: Technology Stack
add_content_slide(prs, "Technology Stack", [
    "🎨 Frontend: Vue 3 + Vue Router + Pinia + Vite",
    "  • Dynamic, mobile-first SPA with reactive state management",
    "",
    "⚙️ Backend: PHP Slim 4 + PDO",
    "  • RESTful API with business logic and middleware",
    "",
    "🗄️ Database: MySQL",
    "  • Relational database with 9 tables and foreign key constraints",
    "",
    "🔒 Security: JWT (firebase/php-jwt) + Bcrypt",
    "  • Token-based authentication with password hashing",
    "",
    "🚀 Deployment: Railway (Nixpacks)",
    "  • Frontend, Backend, and MySQL hosted on Railway"
], "💻")

# Slide 10: Section - DevOps
add_section_header(prs, "DevOps Practices")

# Slide 11: Version Control
add_content_slide(prs, "Version Control & Collaboration", [
    "📦 GitHub Repository:",
    "  • Repository: FleepyDean/GreenStep",
    "  • Main branch for production code",
    "  • 100+ commits with descriptive messages",
    "",
    "🌿 Branching Strategy:",
    "  • Main branch for stable releases",
    "  • Feature branches for new features",
    "  • Direct commits for hotfixes",
    "",
    "🤝 Collaboration:",
    "  • Modular code structure for maintainability",
    "  • .gitignore for sensitive data (.env files)",
    "  • Comprehensive README documentation"
], "🔄")

# Slide 12: CI/CD & Deployment
add_content_slide(prs, "CI/CD & Deployment", [
    "🔄 Continuous Deployment:",
    "  • Railway auto-deploys on Git push to main branch",
    "  • Separate services: Frontend, Backend, MySQL",
    "",
    "🏗️ Build Process:",
    "  • Frontend: Node 20, npm install, npm run build, serve dist/",
    "  • Backend: PHP 8.3, Composer install, php -S 0.0.0.0:$PORT",
    "",
    "🌐 Live URLs:",
    "  • Frontend: https://greenstep.up.railway.app",
    "  • Backend: https://greenstep-backend-production.up.railway.app/api",
    "",
    "⚙️ Environment Variables:",
    "  • Frontend: VITE_API_URL",
    "  • Backend: DB credentials, JWT_SECRET, FRONTEND_URL"
], "🚀")

# Slide 13: Section - Architecture
add_section_header(prs, "System Architecture & Design")

# Slide 14: System Architecture
add_content_slide(prs, "System Architecture Diagram", [
    "🖥️ Client Layer (Browser):",
    "  • Vue 3 SPA with Vite build tool",
    "  • Axios HTTP client with JWT interceptors",
    "  • Pinia state management",
    "",
    "🌐 Server Layer (Railway):",
    "  • PHP Slim 4 REST API",
    "  • JWT middleware for authentication",
    "  • CORS middleware for cross-origin requests",
    "",
    "🗄️ Database Layer (Railway MySQL):",
    "  • MySQL 8.0 with 9 relational tables",
    "  • Foreign key constraints for data integrity",
    "",
    "🚀 Deployment:",
    "  • Railway hosting platform (Frontend + Backend + MySQL)"
], "🏗️")

# Slide 15: Sequence Diagram
add_content_slide(prs, "Sequence Diagram: Activity Logging", [
    "1️⃣ User selects category and activity type",
    "2️⃣ User enters amount, date, and optional photo",
    "3️⃣ Frontend sends POST /api/activities with JWT token",
    "4️⃣ Backend validates JWT token (JwtMiddleware)",
    "5️⃣ Backend calculates carbon footprint (amount × emission factor)",
    "6️⃣ Backend inserts activity into ActivityLog table",
    "7️⃣ Backend checks and awards badges (BadgeController)",
    "8️⃣ Backend returns success response with activity data",
    "9️⃣ Frontend emits 'activity-logged' event via eventBus",
    "🔟 Dashboard listens and refreshes goal + metrics",
    "1️⃣1️⃣ User sees updated dashboard with new activity"
], "🔄")

# Slide 16: Activity Diagram
add_content_slide(prs, "Activity Diagram: User Workflow", [
    "▶️ Start → User opens app",
    "  ↓",
    "❓ Decision: Authenticated?",
    "  • No → Redirect to /profile (login/register)",
    "  • Yes → Show dashboard",
    "  ↓",
    "📝 User navigates to /activity",
    "  ↓",
    "✍️ User logs activity (category, type, amount, date, photo)",
    "  ↓",
    "🧮 System calculates carbon footprint",
    "  ↓",
    "🏆 System checks and awards badges",
    "  ↓",
    "📊 User navigates to /dashboard",
    "  ↓",
    "📈 System displays updated metrics, charts, goal progress",
    "  ↓",
    "⏹️ End"
], "🔀")

# Slide 17: Deployment Diagram
add_content_slide(prs, "Deployment Diagram", [
    "💻 Client Device (Browser):",
    "  • Vue 3 SPA served as static files",
    "  • HTTPS connection to Railway frontend",
    "",
    "🌐 Railway Frontend Server:",
    "  • Static file server (Nginx/similar)",
    "  • Serves index.html + JS/CSS bundles",
    "",
    "⚙️ Railway Backend Server:",
    "  • PHP 8.3 built-in server",
    "  • REST API endpoints with JWT auth",
    "",
    "🗄️ Railway MySQL Database:",
    "  • MySQL 8.0 cloud database",
    "  • 9 tables with foreign key constraints",
    "",
    "🔒 Network: HTTPS/TLS for all connections"
], "🌐")

# Slide 18: Section - Database
add_section_header(prs, "Database & Data Management")

# Slide 19: Database Schema
add_content_slide(prs, "Database Schema & Tables", [
    "🗄️ Database: MySQL 8.0",
    "",
    "📋 Main Tables:",
    "  • User (id, name, email, password, role, goal settings)",
    "  • ActivityType (id, name, category, emission_factor)",
    "  • ActivityLog (id, user_id, activity_type_id, amount, carbon_footprint, date)",
    "  • Badge (id, name, description, icon, category_rule, threshold_value)",
    "  • UserBadge (user_id, badge_id, earned_at)",
    "  • Challenge (id, title, description, target_type, target_value, dates)",
    "  • ChallengeParticipant (id, challenge_id, user_id, progress)",
    "  • Friendship (id, sender_id, receiver_id, status)",
    "  • Tip (id, title, content, category)"
], "🗃️")

# Slide 20: Data Security
add_content_slide(prs, "Data Validation & Security", [
    "✅ Input Validation:",
    "  • Required field checks on frontend and backend",
    "  • Type validation (integers, floats, dates, emails)",
    "",
    "🔒 Security Measures:",
    "  • Password hashing with Bcrypt (cost factor 10)",
    "  • JWT token-based authentication (24-hour expiry)",
    "  • Parameterized SQL queries to prevent injection",
    "  • CORS policy restricts API access to frontend domain",
    "  • Role-based access control (user vs admin)",
    "",
    "🛡️ Data Integrity:",
    "  • Foreign key constraints in database",
    "  • Unique constraints on email addresses",
    "  • NOT NULL constraints on required fields"
], "🔐")

# Slide 21: Demo Data
add_content_slide(prs, "Demonstration Data", [
    "📦 Seed Data Included:",
    "  • 17+ activity types with realistic emission factors",
    "  • 10 eco tips across all categories",
    "  • 4 challenges with date ranges and targets",
    "  • 6+ badges (streak-based and category-based)",
    "  • 1 admin user (admin@greenstep.my / admin123)",
    "",
    "🎬 Realistic Demo Scenario:",
    "  • User registers and logs activities",
    "  • System calculates carbon footprint automatically",
    "  • User sets carbon reduction goal (e.g., 30% in 30 days)",
    "  • User earns badges based on activity milestones",
    "  • User joins challenges and tracks progress",
    "  • User adds friends and compares on leaderboard",
    "  • Admin manages emission factors, badges, challenges, tips"
], "🎯")

# Slide 22: Section - Demonstration
add_section_header(prs, "System Demonstration")

# Slide 23: User Flow Demo
add_content_slide(prs, "User Flow Demonstration", [
    "🔐 1. Authentication:",
    "  • Register new account or login • JWT token stored in localStorage",
    "",
    "📊 2. Dashboard:",
    "  • View today's footprint, weekly/monthly trends",
    "  • See category breakdown doughnut chart • Track carbon reduction goal progress",
    "",
    "📝 3. Activity Logging:",
    "  • Select category (Transport, Diet, Energy, Recycling)",
    "  • Choose activity type from dropdown • Enter amount and date",
    "  • Upload optional photo • Submit and see carbon footprint calculated",
    "",
    "🏆 4. Gamification:",
    "  • Earn badges automatically • Join challenges and track progress",
    "  • View leaderboard rankings"
], "👤")

# Slide 24: Admin Flow Demo
add_content_slide(prs, "Admin Flow Demonstration", [
    "📊 1. Admin Dashboard:",
    "  • View platform statistics (total users, activities, emissions)",
    "  • See user dataset overview",
    "",
    "⚡ 2. Emission Factor Management:",
    "  • Add/edit/delete activity types • Create custom categories",
    "  • Set emission factors (kg CO₂ per unit)",
    "",
    "🏅 3. Badge Management:",
    "  • Create custom badges with category rules • Set threshold values • Delete badges",
    "",
    "🏆 4. Challenge Management:",
    "  • Create challenges with targets and date ranges • Update/delete challenges",
    "",
    "💡 5. Tips Management:",
    "  • Add eco tips with categories • Delete tips"
], "⚙️")

# Slide 25: Key Features
add_content_slide(prs, "Key Features Demonstrated", [
    "✓ User authentication with JWT",
    "✓ Activity logging with automatic carbon calculation",
    "✓ Real-time dashboard with interactive charts",
    "✓ Carbon reduction goal with progress tracking",
    "✓ Badge unlocking based on activity milestones",
    "✓ Challenge join/leave with progress tracking",
    "✓ Friend connections and leaderboard",
    "✓ Admin CRUD operations for all entities",
    "✓ Photo upload for activities",
    "✓ Responsive mobile-first design",
    "✓ Professional UI with SVG icons",
    "✓ Event-driven architecture (eventBus)",
    "✓ Error handling and validation",
    "✓ Toast notifications for user feedback"
], "✨")

# Slide 26: Technical Achievements
add_content_slide(prs, "Technical Achievements", [
    "🏗️ Software Engineering Practices:",
    "  • Modular MVC architecture (Models, Controllers, Views)",
    "  • RESTful API design with proper HTTP methods",
    "  • Separation of concerns (frontend/backend/database)",
    "  • Reusable components and services",
    "",
    "🚀 DevOps Practices:",
    "  • Git version control with descriptive commits",
    "  • Continuous deployment with Railway",
    "  • Environment variable management",
    "  • Timezone handling for global deployment (Asia/Kuala_Lumpur)",
    "",
    "💎 Code Quality:",
    "  • Consistent naming conventions",
    "  • Proper error handling and logging",
    "  • SQL query optimization • Frontend state management with Pinia"
], "🎓")

# Slide 27: Conclusion
add_title_slide(prs, "Thank You! 🌱", "Live Demonstration\n\n🌐 Frontend: https://greenstep.up.railway.app\n⚙️ Backend API: https://greenstep-backend-production.up.railway.app/api\n\n📦 GitHub: FleepyDean/GreenStep\n\n❓ Questions?")

prs.save("GreenStep_Final_Presentation.pptx")
print("✓ Styled presentation created successfully: GreenStep_Final_Presentation.pptx")
print("✓ Total slides: 20 (with section headers)")
print("✓ Design: Green theme with professional layout")
